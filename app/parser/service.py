"""
Document parsing: multi-format to structured Markdown/JSON.
Primary engine: Docling (preserves tables, headings, OCR for scanned PDFs).
Fallback: legacy parsers (PyMuPDF, python-docx, openpyxl, python-pptx).
"""

import io
import logging
import tempfile
from pathlib import Path

from app.core.config import settings
from app.models.parser import ParsedDocument, ParsedDocumentMetadata

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Docling-based parser (primary)
# ---------------------------------------------------------------------------


def _parse_with_docling(content: bytes, filename: str) -> ParsedDocument:
    """
    Convert document to structured Markdown using Docling.
    Supports PDF, DOCX, XLSX, PPTX with table/heading preservation.
    """
    from docling.document_converter import DocumentConverter

    suffix = Path(filename).suffix.lower()

    # Docling requires a file path — write to temp file
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        converter = DocumentConverter()
        result = converter.convert(tmp_path)

        # Export structured Markdown (preserves tables, headings, lists)
        markdown_content = result.document.export_to_markdown()

        # Export JSON structure for metadata/raw access
        raw_structure = result.document.export_to_dict()

        if not markdown_content or not markdown_content.strip():
            markdown_content = "(no text extracted by Docling)"

        return ParsedDocument(
            content=markdown_content.strip(),
            raw_structure=raw_structure,
            metadata=ParsedDocumentMetadata(
                filename=_safe_filename(filename),
                type=suffix.lstrip("."),
                parser_engine="docling",
            ),
        )
    finally:
        Path(tmp_path).unlink(missing_ok=True)


# ---------------------------------------------------------------------------
# Legacy parsers (fallback)
# ---------------------------------------------------------------------------


def _safe_filename(name: str) -> str:
    """Sanitise filename for metadata."""
    return Path(name).name[:255] if name else "unknown"


def _parse_pdf(content: bytes, filename: str) -> ParsedDocument:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=content, filetype="pdf")
    parts = []
    for page in doc:
        parts.append(page.get_text())
    doc.close()
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type="pdf", parser_engine="legacy"
        ),
    )


def _parse_docx(content: bytes, filename: str) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type="docx", parser_engine="legacy"
        ),
    )


def _parse_xlsx(content: bytes, filename: str) -> ParsedDocument:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## {sheet.title}\n")
        for row in sheet.iter_rows(values_only=True):
            parts.append(" | ".join(str(c) if c is not None else "" for c in row))
    wb.close()
    text = "\n".join(parts).strip() or "(no content)"
    return ParsedDocument(
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type="xlsx", parser_engine="legacy"
        ),
    )


def _parse_pptx(content: bytes, filename: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"## Slide {i + 1}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type="pptx", parser_engine="legacy"
        ),
    )


def _parse_plain(content: bytes, filename: str, content_type: str) -> ParsedDocument:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1", errors="replace")
    return ParsedDocument(
        content=text.strip() or "(empty)",
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename),
            type=content_type,
            parser_engine="legacy",
        ),
    )


def _parse_mermaid(content: bytes, filename: str, content_type: str) -> ParsedDocument:
    doc = _parse_plain(content, filename, content_type)
    doc.content = f"```mermaid\n{doc.content}\n```"
    return doc


# ---------------------------------------------------------------------------
# Extension whitelist and dispatch
# ---------------------------------------------------------------------------

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".txt",
    ".md",
    ".mmd",
    ".mermaid",
}

# Extensions that Docling can handle natively
_DOCLING_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}

_LEGACY_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".txt": lambda c, f: _parse_plain(c, f, "txt"),
    ".md": lambda c, f: _parse_plain(c, f, "md"),
    ".mmd": lambda c, f: _parse_mermaid(c, f, "mmd"),
    ".mermaid": lambda c, f: _parse_mermaid(c, f, "mermaid"),
}


def parse_file(content: bytes, filename: str) -> ParsedDocument:
    """
    Parse a single file into ParsedDocument.

    Engine selection (via PARSER_ENGINE config):
      - "auto": Try Docling first for supported formats, fall back to legacy.
      - "docling": Use Docling only (raises on failure).
      - "legacy": Use legacy parsers only (original behavior).

    Raises ValueError for unsupported file types.
    """
    path = Path(filename)
    suffix = path.suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {suffix}. Allowed: {sorted(ALLOWED_EXTENSIONS)}"
        )

    engine = settings.PARSER_ENGINE

    # Text-like formats always use legacy (no benefit from Docling)
    if suffix in {".txt", ".md", ".mmd", ".mermaid"}:
        return _LEGACY_PARSERS[suffix](content, filename)

    # Docling-capable formats
    if engine == "docling":
        return _parse_with_docling(content, filename)

    if engine == "auto" and suffix in _DOCLING_EXTENSIONS:
        try:
            return _parse_with_docling(content, filename)
        except Exception as e:
            logger.warning(
                "Docling failed for %s, falling back to legacy parser: %s",
                filename,
                e,
            )
            return _LEGACY_PARSERS[suffix](content, filename)

    # engine == "legacy" or fallback
    return _LEGACY_PARSERS[suffix](content, filename)
